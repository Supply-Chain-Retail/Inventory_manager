import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Dark Theme)
    BG_COLOR = RGBColor(15, 23, 42)       # Slate 900
    ACCENT_TEAL = RGBColor(20, 184, 166)  # Teal 500
    ACCENT_SKY = RGBColor(14, 165, 233)   # Sky 500
    TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="OMNIOPT AI"):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Calibri"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_TEAL

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Accent color block on left
    shape = slide1.shapes.add_shape(
        1,  # Rectangle
        Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_TEAL
    shape.line.fill.background()

    # Title & Subtitle in single textbox
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11), Inches(3.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0

    p_title = tf1.paragraphs[0]
    p_title.text = "OmniOpt AI"
    p_title.font.name = "Calibri"
    p_title.font.size = Pt(64)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    p_subtitle = tf1.add_paragraph()
    p_subtitle.text = "Predictive Logistics & Fulfillment Platform"
    p_subtitle.font.name = "Calibri"
    p_subtitle.font.size = Pt(28)
    p_subtitle.font.color.rgb = ACCENT_SKY
    p_subtitle.space_before = Pt(10)

    p_brief = tf1.add_paragraph()
    p_brief.text = "Staff Operations Guide & System Architecture"
    p_brief.font.name = "Calibri"
    p_brief.font.size = Pt(16)
    p_brief.font.color.rgb = TEXT_MUTED
    p_brief.space_before = Pt(40)

    # -------------------------------------------------------------
    # SLIDE 2: The E-Commerce Challenge
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Logistics Challenge We Solve", "Market Pain Points")

    # 4 Cards/Columns for challenges
    cols = [
        {"title": "Stockout vs Overstock", "desc": "Imbalances lead to lost sales or high storage costs. Manual restocks are slow and error-prone."},
        {"title": "Static Reordering", "desc": "Traditional inventory management uses static reorder points that do not adapt to fast-changing demand velocity."},
        {"title": "Environmental Factors", "desc": "Weather shifts (heatwaves, heavy rain) and regional trends drastically shift demand overnight without warning."},
        {"title": "Fulfillment Friction", "desc": "Inefficient inter-hub shipping routes and delayed factory restock coordination degrade delivery speeds."}
    ]

    left_margin = Inches(0.8)
    card_width = Inches(2.65)
    card_gap = Inches(0.35)
    top_pos = Inches(2.2)
    card_height = Inches(4.2)

    for i, col in enumerate(cols):
        l_pos = left_margin + i * (card_width + card_gap)
        
        # Draw background card shape
        card = slide2.shapes.add_shape(1, l_pos, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_TEAL if i == 0 else CARD_BG
        
        # Add Text Frame inside the card
        tb = slide2.shapes.add_textbox(l_pos + Inches(0.2), top_pos + Inches(0.3), card_width - Inches(0.4), card_height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Number indicator
        p_num = tf.paragraphs[0]
        p_num.text = f"0{i+1}"
        p_num.font.name = "Calibri"
        p_num.font.size = Pt(20)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_TEAL
        
        # Title
        p_title = tf.add_paragraph()
        p_title.text = col["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.space_before = Pt(12)
        p_title.space_after = Pt(8)
        
        # Description
        p_desc = tf.add_paragraph()
        p_desc.text = col["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "System Architecture & Flow", "How It Works")

    # Draw diagram boxes (Horizontal Flow: Data -> Prediction -> Decision -> Fulfillment)
    flow_steps = [
        {"title": "1. Live Sales Feed", "subtitle": "WebSockets Simulator", "desc": "Streams real-time sales from Amazon, Myntra, Flipkart, Meesho every 3 seconds."},
        {"title": "2. Predictive Brain", "subtitle": "Hybrid ML Ensemble", "desc": "Forecasts regional demand utilizing Scikit-learn HGB and Random Forest regressors."},
        {"title": "3. Decision Engine", "subtitle": "Inventory Logic", "desc": "Auto-detects stockouts; rebalances from surplus hubs or requests factory orders."},
        {"title": "4. Fulfillment Map", "subtitle": "Routing Engine", "desc": "Visualizes center coordinates, dispatch paths, restock countdowns via Leaflet.js."}
    ]

    card_width = Inches(2.7)
    card_gap = Inches(0.3)
    left_margin = Inches(0.8)
    top_pos = Inches(2.2)
    card_height = Inches(4.0)

    for i, step in enumerate(flow_steps):
        l_pos = left_margin + i * (card_width + card_gap)
        
        # Background card
        card = slide3.shapes.add_shape(1, l_pos, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_SKY
        
        tb = slide3.shapes.add_textbox(l_pos + Inches(0.2), top_pos + Inches(0.3), card_width - Inches(0.4), card_height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_title = tf.paragraphs[0]
        p_title.text = step["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        
        p_sub = tf.add_paragraph()
        p_sub.text = step["subtitle"]
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ACCENT_TEAL
        p_sub.space_before = Pt(4)
        p_sub.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = step["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

        # Draw connecting arrow to the next slide card (except last card)
        if i < len(flow_steps) - 1:
            arrow_l = l_pos + card_width + Inches(0.05)
            arrow_w = card_gap - Inches(0.1)
            arrow = slide3.shapes.add_shape(
                34,  # Right Arrow
                arrow_l, top_pos + Inches(1.8), arrow_w, Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_SKY
            arrow.line.fill.background()

    # -------------------------------------------------------------
    # SLIDE 4: Machine Learning Forecast Brain
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Intelligent Ensemble Demand Forecasting", "Machine Learning Brain")

    # Left Column: Model Specs & Accuracy
    tb_left = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

    p_model_title = tf_left.paragraphs[0]
    p_model_title.text = "Hybrid Ensemble Model"
    p_model_title.font.name = "Calibri"
    p_model_title.font.size = Pt(22)
    p_model_title.font.bold = True
    p_model_title.font.color.rgb = ACCENT_TEAL

    p_model_desc = tf_left.add_paragraph()
    p_model_desc.text = "We combine two state-of-the-art regressors to balance local precision and global generalization:"
    p_model_desc.font.name = "Calibri"
    p_model_desc.font.size = Pt(14)
    p_model_desc.font.color.rgb = TEXT_MUTED
    p_model_desc.space_before = Pt(8)
    p_model_desc.space_after = Pt(14)

    bullets = [
        ("HistGradientBoosting (60% Weight):", " Handles large-scale numeric datasets and non-linear patterns at lightning speed."),
        ("RandomForest (40% Weight):", " Prevents overfitting, ensuring reliable projections even with erratic sales peaks."),
        ("Advanced Feature Engineering:", " Models use historical baselines, monthly/weekly seasonality, and active store prices to predict sales.")
    ]
    for title, desc in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = "• "
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(6)
        
        run_title = p_b.add_run()
        run_title.text = title
        run_title.font.bold = True
        run_title.font.color.rgb = TEXT_WHITE
        
        run_desc = p_b.add_run()
        run_desc.text = desc
        run_desc.font.color.rgb = TEXT_MUTED

    # Right Column: Big Stat Box
    stat_card = slide4.shapes.add_shape(1, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    stat_card.fill.solid()
    stat_card.fill.fore_color.rgb = CARD_BG
    stat_card.line.color.rgb = ACCENT_TEAL
    stat_card.line.width = Pt(1.5)

    tb_stat = slide4.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(4.9), Inches(3.9))
    tf_stat = tb_stat.text_frame
    tf_stat.word_wrap = True
    tf_stat.margin_left = tf_stat.margin_right = tf_stat.margin_top = tf_stat.margin_bottom = 0

    p_stat_title = tf_stat.paragraphs[0]
    p_stat_title.alignment = PP_ALIGN.CENTER
    p_stat_title.text = "PREDICTION ACCURACY"
    p_stat_title.font.name = "Calibri"
    p_stat_title.font.size = Pt(14)
    p_stat_title.font.bold = True
    p_stat_title.font.color.rgb = TEXT_MUTED

    p_stat_num = tf_stat.add_paragraph()
    p_stat_num.alignment = PP_ALIGN.CENTER
    p_stat_num.text = "94.6%"
    p_stat_num.font.name = "Calibri"
    p_stat_num.font.size = Pt(88)
    p_stat_num.font.bold = True
    p_stat_num.font.color.rgb = ACCENT_TEAL
    p_stat_num.space_before = Pt(0)

    p_stat_sub = tf_stat.add_paragraph()
    p_stat_sub.alignment = PP_ALIGN.CENTER
    p_stat_sub.text = "R² Coefficient of Determination"
    p_stat_sub.font.name = "Calibri"
    p_stat_sub.font.size = Pt(16)
    p_stat_sub.font.bold = True
    p_stat_sub.font.color.rgb = TEXT_WHITE
    p_stat_sub.space_before = Pt(0)

    p_stat_footer = tf_stat.add_paragraph()
    p_stat_footer.alignment = PP_ALIGN.CENTER
    p_stat_footer.text = "Calculated across real-time multi-platform transaction logs. Retrained continuously on live sales."
    p_stat_footer.font.name = "Calibri"
    p_stat_footer.font.size = Pt(12)
    p_stat_footer.font.color.rgb = TEXT_MUTED
    p_stat_footer.space_before = Pt(20)

    # -------------------------------------------------------------
    # SLIDE 5: Real-time Rebalancing & Factory Replenishment
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Automated Rebalancing & Factory restocks", "Replenishment Engine")

    # Left Column: Rebalancing
    card_l = slide5.shapes.add_shape(1, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = CARD_BG

    tb_cl = slide5.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(5.0), Inches(3.9))
    tf_cl = tb_cl.text_frame
    tf_cl.word_wrap = True
    tf_cl.margin_left = tf_cl.margin_right = tf_cl.margin_top = tf_cl.margin_bottom = 0

    p_cl_title = tf_cl.paragraphs[0]
    p_cl_title.text = "🔄 Inter-Hub Rebalancing"
    p_cl_title.font.name = "Calibri"
    p_cl_title.font.size = Pt(20)
    p_cl_title.font.bold = True
    p_cl_title.font.color.rgb = ACCENT_TEAL
    p_cl_title.space_after = Pt(14)

    bullets_l = [
        "Identifies stock deficits at any hub automatically.",
        "Scans geographic regions for nearby hubs with stock surplus.",
        "Initiates direct inter-hub transfer (e.g. North Hub to West Hub).",
        "Avoids redundant warehousing costs and minimizes transport time."
    ]
    for b in bullets_l:
        p_b = tf_cl.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(10)

    # Right Column: Factory Restock
    card_r = slide5.shapes.add_shape(1, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = CARD_BG

    tb_cr = slide5.shapes.add_textbox(Inches(7.2), Inches(2.3), Inches(5.0), Inches(3.9))
    tf_cr = tb_cr.text_frame
    tf_cr.word_wrap = True
    tf_cr.margin_left = tf_cr.margin_right = tf_cr.margin_top = tf_cr.margin_bottom = 0

    p_cr_title = tf_cr.paragraphs[0]
    p_cr_title.text = "🏭 Direct Factory Restocking"
    p_cr_title.font.name = "Calibri"
    p_cr_title.font.size = Pt(20)
    p_cr_title.font.bold = True
    p_cr_title.font.color.rgb = ACCENT_SKY
    p_cr_title.space_after = Pt(14)

    bullets_r = [
        "Triggered when no nearby hubs have surplus stock.",
        "Generates automated Factory Restock Order directly.",
        "Real-time countdown tracker displays delivery time.",
        "Physical inventory is updated and locked only upon countdown arrival to ensure strict consistency."
    ]
    for b in bullets_r:
        p_b = tf_cr.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 6: Weather Telemetry Integration
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Weather-Driven Dynamic Safety Stock", "Environmental Telemetry")

    # Left TextBox
    tb_w_left = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_w_l = tb_w_left.text_frame
    tf_w_l.word_wrap = True
    tf_w_l.margin_left = tf_w_l.margin_right = tf_w_l.margin_top = tf_w_l.margin_bottom = 0

    p_w_title = tf_w_l.paragraphs[0]
    p_w_title.text = "Open-Meteo API Integration"
    p_w_title.font.name = "Calibri"
    p_w_title.font.size = Pt(22)
    p_w_title.font.bold = True
    p_w_title.font.color.rgb = ACCENT_SKY

    p_w_desc = tf_w_l.add_paragraph()
    p_w_desc.text = "OmniOpt AI integrates real-time climate telemetry directly into forecasting, adjusting safety stock margins on-the-fly based on atmospheric conditions."
    p_w_desc.font.name = "Calibri"
    p_w_desc.font.size = Pt(14)
    p_w_desc.font.color.rgb = TEXT_MUTED
    p_w_desc.space_before = Pt(8)
    p_w_desc.space_after = Pt(14)

    bullets_w = [
        ("Rainfall Multiplier:", " Increases safety buffers for waterproof gear, umbrellas, and select home-delivery goods by 1.25x."),
        ("Temperature Multiplier:", " Heavy summer heatwaves trigger automatic 1.30x scaling for grocery, beverage, and cooling products."),
        ("Risk Mitigation:", " System automatically locks in higher safety stock during storms/extreme weather events to protect against road delays.")
    ]
    for title, desc in bullets_w:
        p_b = tf_w_l.add_paragraph()
        p_b.text = "• "
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(6)
        
        run_title = p_b.add_run()
        run_title.text = title
        run_title.font.bold = True
        run_title.font.color.rgb = TEXT_WHITE
        
        run_desc = p_b.add_run()
        run_desc.text = desc
        run_desc.font.color.rgb = TEXT_MUTED

    # Right layout: 3 horizontal stats cards showing multipliers
    multipliers = [
        {"title": "Rain / Storms", "mult": "1.25x", "cat": "Clothing & Delivery Items"},
        {"title": "Extreme Heat", "mult": "1.30x", "cat": "Groceries & Beverages"},
        {"title": "Default Climate", "mult": "1.00x", "cat": "Baseline Demand Profile"}
    ]

    r_left = Inches(7.5)
    r_width = Inches(5.0)
    card_h = Inches(1.3)
    card_gap = Inches(0.2)

    for i, m in enumerate(multipliers):
        r_top = Inches(2.0) + i * (card_h + card_gap)
        
        m_card = slide6.shapes.add_shape(1, r_left, r_top, r_width, card_h)
        m_card.fill.solid()
        m_card.fill.fore_color.rgb = CARD_BG
        m_card.line.color.rgb = ACCENT_TEAL if i == 1 else CARD_BG
        
        tb_m = slide6.shapes.add_textbox(r_left + Inches(0.3), r_top + Inches(0.2), r_width - Inches(0.6), card_h - Inches(0.4))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = tf_m.margin_right = tf_m.margin_top = tf_m.margin_bottom = 0
        
        # Multiplier on the right, label on the left (split using single layout)
        p_m = tf_m.paragraphs[0]
        p_m.text = m["title"] + "  |  "
        p_m.font.name = "Calibri"
        p_m.font.size = Pt(16)
        p_m.font.bold = True
        p_m.font.color.rgb = TEXT_WHITE
        
        run_mult = p_m.add_run()
        run_mult.text = m["mult"]
        run_mult.font.bold = True
        run_mult.font.color.rgb = ACCENT_TEAL
        
        p_m_sub = tf_m.add_paragraph()
        p_m_sub.text = m["cat"]
        p_m_sub.font.name = "Calibri"
        p_m_sub.font.size = Pt(12)
        p_m_sub.font.color.rgb = TEXT_MUTED
        p_m_sub.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 7: Staff Dashboard Interface Guide
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Operational Staff Dashboard Guide", "Control Center")

    # Grid of 4 features
    features = [
        {"title": "🗺️ Geographic Route Engine", "desc": "Interactive Leaflet.js map. Shows fulfillment hubs (North, South, East, West) and displays route geometry overlay during active stock dispatches."},
        {"title": "📊 Real-Time Forecasting Grid", "desc": "Displays tomorrow's projected demand against current inventory. Highlights under-stocked categories in bright warning orange."},
        {"title": "🔔 Active Transit Panel", "desc": "Tracks inbound restocks and transfers. Real-time countdown clock ticks down. Prevents staff from double-ordering identical stock items."},
        {"title": "📈 Live E-Commerce Simulator", "desc": "Simulates continuous online purchases (Amazon, Flipkart, etc.) decrementing physical inventory. Retrains ML engine instantly."}
    ]

    col_w = Inches(5.6)
    row_h = Inches(2.0)
    left_m = Inches(0.8)
    top_m = Inches(2.0)
    gap_x = Inches(0.4)
    gap_y = Inches(0.4)

    for i, feat in enumerate(features):
        col_idx = i % 2
        row_idx = i // 2
        
        cx = left_m + col_idx * (col_w + gap_x)
        cy = top_m + row_idx * (row_h + gap_y)
        
        card = slide7.shapes.add_shape(1, cx, cy, col_w, row_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        
        tb = slide7.shapes.add_textbox(cx + Inches(0.3), cy + Inches(0.2), col_w - Inches(0.6), row_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_t = tf.paragraphs[0]
        p_t.text = feat["title"]
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_after = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = feat["desc"]
        p_d.font.name = "Calibri"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 8: Deployment & Technical Stack
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Platform Architecture & Deployments", "Deployment Blueprint")

    # Left: Frontend
    f_card = slide8.shapes.add_shape(1, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5))
    f_card.fill.solid()
    f_card.fill.fore_color.rgb = CARD_BG
    f_card.line.color.rgb = ACCENT_TEAL

    tb_f = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(3.2), Inches(4.1))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0

    p_f_t = tf_f.paragraphs[0]
    p_f_t.text = "Frontend Client"
    p_f_t.font.name = "Calibri"
    p_f_t.font.size = Pt(18)
    p_f_t.font.bold = True
    p_f_t.font.color.rgb = ACCENT_TEAL
    p_f_t.space_after = Pt(12)

    bullets_f = [
        "Tech: HTML5, TailwindCSS, Chart.js, Leaflet.js mapping.",
        "Production: Hosted statically on GitHub Pages.",
        "Fallback mode: Connects to local JSON data model dynamically if backend API drops."
    ]
    for b in bullets_f:
        p_b = tf_f.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(8)

    # Middle: Backend
    b_card = slide8.shapes.add_shape(1, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.5))
    b_card.fill.solid()
    b_card.fill.fore_color.rgb = CARD_BG
    b_card.line.color.rgb = ACCENT_SKY

    tb_b = slide8.shapes.add_textbox(Inches(5.0), Inches(2.2), Inches(3.2), Inches(4.1))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0

    p_b_t = tf_b.paragraphs[0]
    p_b_t.text = "Backend Server"
    p_b_t.font.name = "Calibri"
    p_b_t.font.size = Pt(18)
    p_b_t.font.bold = True
    p_b_t.font.color.rgb = ACCENT_SKY
    p_b_t.space_after = Pt(12)

    bullets_b = [
        "Tech: FastAPI framework, Scikit-learn regressors, Uvicorn, Pandas.",
        "Production: Hosted on Render cloud container service.",
        "Real-Time Sync: Employs persistent WebSockets stream for state broadcasts."
    ]
    for b in bullets_b:
        p_b = tf_b.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(8)

    # Right: Data
    d_card = slide8.shapes.add_shape(1, Inches(8.8), Inches(2.0), Inches(3.7), Inches(4.5))
    d_card.fill.solid()
    d_card.fill.fore_color.rgb = CARD_BG
    d_card.line.color.rgb = TEXT_MUTED

    tb_d = slide8.shapes.add_textbox(Inches(9.0), Inches(2.2), Inches(3.3), Inches(4.1))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_right = tf_d.margin_top = tf_d.margin_bottom = 0

    p_d_t = tf_d.paragraphs[0]
    p_d_t.text = "Database Model"
    p_d_t.font.name = "Calibri"
    p_d_t.font.size = Pt(18)
    p_d_t.font.bold = True
    p_d_t.font.color.rgb = TEXT_WHITE
    p_d_t.space_after = Pt(12)

    bullets_d = [
        "Model: Portable JSON structure (inventory_db.json).",
        "Backup: Auto-commits updates dynamically on action/sales updates.",
        "Seeding: Auto-generates realistic sales history on boot if missing."
    ]
    for b in bullets_d:
        p_b = tf_d.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(8)

    prs.save("/Users/kalliyany/Documents/finale/OmniOpt_AI_Presentation.pptx")
    print("✅ PowerPoint presentation generated successfully at /Users/kalliyany/Documents/finale/OmniOpt_AI_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
